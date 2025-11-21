#!/usr/bin/env python3
"""
Helper script to create showcase content from generated presentations.

This script helps you:
1. Extract key slides from presentations
2. Create screenshots for documentation
3. Generate showcase gallery entries

Usage:
    python scripts/showcase/create_showcase.py --presentation path/to/presentation.pptx
"""

import argparse
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:
    print("Error: python-pptx not installed. Run: pip install -r requirements.txt")
    sys.exit(1)


def extract_slides_info(presentation_path, output_dir):
    """Extract information about slides for showcase."""
    prs = Presentation(presentation_path)
    
    print(f"📊 Presentation: {presentation_path}")
    print(f"   Total slides: {len(prs.slides)}")
    print(f"   Output directory: {output_dir}")
    
    # Create output directory
    output_dir = Path(output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)
    
    # Generate slide summary
    summary = []
    summary.append(f"# Presentation Summary\n")
    summary.append(f"**File**: `{presentation_path.name}`\n")
    summary.append(f"**Total Slides**: {len(prs.slides)}\n\n")
    summary.append("## Key Slides\n\n")
    
    for i, slide in enumerate(prs.slides, 1):
        # Get slide title if available
        title = "Untitled"
        if slide.shapes.title:
            title = slide.shapes.title.text[:50]  # Truncate long titles
        
        summary.append(f"{i}. **{title}**\n")
    
    # Save summary
    summary_path = output_dir / "presentation_summary.md"
    summary_path.write_text("".join(summary))
    print(f"✅ Summary saved to: {summary_path}")
    
    print("\n💡 Next steps:")
    print("   1. Open the presentation and take screenshots of key slides")
    print("   2. Save screenshots to showcase/screenshots/")
    print("   3. Update showcase/README.md with descriptions")
    print("   4. Add images to README.md in the Showcase section")


def main():
    parser = argparse.ArgumentParser(
        description="Create showcase content from presentations"
    )
    parser.add_argument(
        "--presentation",
        type=str,
        required=True,
        help="Path to the presentation file (.pptx)"
    )
    parser.add_argument(
        "--output",
        type=str,
        default="showcase/thales-example",
        help="Output directory for showcase content"
    )
    
    args = parser.parse_args()
    
    presentation_path = Path(args.presentation)
    if not presentation_path.exists():
        print(f"❌ Error: Presentation not found: {presentation_path}")
        sys.exit(1)
    
    extract_slides_info(presentation_path, args.output)


if __name__ == "__main__":
    main()

