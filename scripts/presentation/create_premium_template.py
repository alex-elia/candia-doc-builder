#!/usr/bin/env python3
"""
Create a premium tech/sustainable template for PowerPoint
This creates a .potx template file that can be reused
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Premium Tech Sustainable Color Palette
# Inspired by sustainable tech companies (Tesla, Apple, etc.)
DARK_GREEN = RGBColor(0, 100, 80)      # Sustainable green
TECH_BLUE = RGBColor(0, 120, 215)      # Modern tech blue
LIGHT_GREEN = RGBColor(144, 238, 144)  # Light sustainable green
DARK_GRAY = RGBColor(30, 30, 30)       # Dark background
LIGHT_GRAY = RGBColor(240, 240, 240)   # Light background
WHITE = RGBColor(255, 255, 255)
ACCENT_GOLD = RGBColor(255, 215, 0)    # Premium accent

def create_premium_template():
    """Create a premium tech/sustainable PowerPoint template"""
    
    prs = Presentation()
    
    # Set 16:9 aspect ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Define custom theme colors (if supported)
    # Note: python-pptx doesn't fully support theme colors, but we'll use RGB
    
    # Create master slide layouts
    # Layout 0: Title slide
    title_slide = prs.slide_layouts[0]
    
    # Layout 1: Title and content
    content_slide = prs.slide_layouts[1]
    
    # Create a sample slide to show template style
    slide = prs.slides.add_slide(title_slide)
    
    # Add background gradient effect (simulated with shapes)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                                Inches(0), Inches(0),
                                prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_GRAY
    bg.line.fill.background()
    bg.z_order = 0  # Send to back
    
    # Add accent line at top
    accent_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                         Inches(0), Inches(0),
                                         prs.slide_width, Inches(0.15))
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = TECH_BLUE
    accent_line.line.fill.background()
    
    # Title
    title = slide.shapes.title
    title.text = "Premium Tech Template"
    title_frame = title.text_frame
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WHITE
    
    # Subtitle
    subtitle = slide.placeholders[1]
    subtitle.text = "Sustainable Technology Presentation"
    subtitle_frame = subtitle.text_frame
    subtitle_frame.paragraphs[0].font.size = Pt(20)
    subtitle_frame.paragraphs[0].font.color.rgb = LIGHT_GREEN
    
    # Save as template
    template_file = "premium_tech_template.potx"
    prs.save(template_file)
    print(f"✓ Template créé: {template_file}")
    
    return template_file

if __name__ == "__main__":
    create_premium_template()

