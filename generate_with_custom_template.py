#!/usr/bin/env python3
"""
Generate Thales presentation using custom template: Modele_Decalage.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

def add_image_safe(slide, image_path, left, top, width=None, height=None):
    """Add image if exists"""
    if os.path.exists(image_path):
        if width and height:
            return slide.shapes.add_picture(image_path, left, top, width, height)
        else:
            return slide.shapes.add_picture(image_path, left, top)
    return None

def create_presentation_with_template(template_path):
    """Create presentation using custom template"""
    
    if not os.path.exists(template_path):
        print(f"❌ Template not found: {template_path}")
        return None
    
    print(f"📄 Loading template: {template_path}")
    prs = Presentation(template_path)
    
    images_dir = "images"
    
    # Get slide dimensions from template
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    print(f"📐 Slide size: {slide_width/Inches(1):.1f}\" x {slide_height/Inches(1):.1f}\"")
    
    # ============================================
    # SLIDE 1: TITLE
    # ============================================
    # Use template's title slide layout if available
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        if slide.shapes.title:
            slide.shapes.title.text = "IA Embarquée pour\nSystèmes Satellitaires"
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = "Une Approche Architecturale pour l'Intelligence Autonome à Bord\n\nAlexandre GON\nArchitecture & Solutions IA"
    except:
        # Fallback: blank slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.3), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = "IA Embarquée pour\nSystèmes Satellitaires"
        title_frame.paragraphs[0].font.size = Pt(48)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add background image if available
    bg_image = os.path.join(images_dir, "background_space.jpg")
    if os.path.exists(bg_image):
        # Add as background (behind existing content)
        bg = slide.shapes.add_picture(bg_image, Inches(0), Inches(0),
                                     slide_width, slide_height)
        # Move to back
        bg.element.getparent().remove(bg.element)
        slide.shapes._spTree.insert(2, bg.element)
    
    # ============================================
    # SLIDE 2: THE CHALLENGE
    # ============================================
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content
        if slide.shapes.title:
            slide.shapes.title.text = "Le Défi"
    except:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "Le Défi"
        title_frame.paragraphs[0].font.size = Pt(48)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Content
    challenges = [
        ("50-500GB/jour", "Transmission coûteuse"),
        ("Délais critiques", "Alertes retardées"),
        ("Pannes communication", "Pas de décision autonome"),
        ("150K-1.8M$/an", "Coûts opérationnels")
    ]
    
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.4), Inches(4.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    for metric, desc in challenges:
        p = content_frame.add_paragraph()
        p.text = f"• {metric}: {desc}"
        p.font.size = Pt(20)
        p.space_after = Pt(12)
        if metric.startswith("150K"):
            p.font.bold = True
    
    # ============================================
    # SLIDE 3: SOLUTION
    # ============================================
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        if slide.shapes.title:
            slide.shapes.title.text = "La Solution"
    except:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "La Solution"
        title_frame.paragraphs[0].font.size = Pt(48)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Value props
    value_props = [
        ("Opérations Autonomes", "80%", "Décisions sans intervention sol"),
        ("Réduction Coûts", "150K-1.8M$/an", "Transmission -80-99%"),
        ("Maintenance Prédictive", "Lifetime +", "Détection précoce"),
        ("Solution Universelle", "Portfolio-wide", "Tous types satellites")
    ]
    
    y_start = 1.8
    for i, (title, metric, desc) in enumerate(value_props):
        y_pos = y_start + i * 1.1
        
        # Card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(1), Inches(y_pos),
                                     Inches(11), Inches(0.9))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(70, 130, 180)
        card.line.color.rgb = RGBColor(15, 32, 66)
        card.line.width = Pt(2)
        
        # Title + Metric
        text_box = slide.shapes.add_textbox(Inches(1.2), Inches(y_pos + 0.1), Inches(8), Inches(0.4))
        text_frame = text_box.text_frame
        text_frame.text = f"{title} - {metric}"
        text_frame.paragraphs[0].font.size = Pt(20)
        text_frame.paragraphs[0].font.bold = True
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(1.2), Inches(y_pos + 0.55), Inches(10), Inches(0.3))
        desc_frame = desc_box.text_frame
        desc_frame.text = desc
        desc_frame.paragraphs[0].font.size = Pt(14)
    
    # ============================================
    # SLIDE 4: ARCHITECTURE
    # ============================================
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        if slide.shapes.title:
            slide.shapes.title.text = "Architecture Système"
    except:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "Architecture Système"
        title_frame.paragraphs[0].font.size = Pt(40)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Try architecture diagram
    arch_image = os.path.join(images_dir, "architecture_diagram.png")
    if os.path.exists(arch_image):
        slide.shapes.add_picture(arch_image, Inches(1.5), Inches(1.5),
                                width=Inches(10), height=Inches(5))
    else:
        # Fallback diagram
        earth_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(2), Inches(1.5),
                                        Inches(9), Inches(1.2))
        earth_box.fill.solid()
        earth_box.fill.fore_color.rgb = RGBColor(70, 130, 180)
        earth_box.line.color.rgb = RGBColor(15, 32, 66)
        earth_box.line.width = Pt(3)
        
        earth_text = slide.shapes.add_textbox(Inches(2.2), Inches(1.7), Inches(8.6), Inches(0.8))
        earth_frame = earth_text.text_frame
        earth_frame.text = "TERRE - Station Terrestre | Backend Central"
        earth_frame.paragraphs[0].font.size = Pt(18)
        earth_frame.paragraphs[0].font.bold = True
        earth_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Arrow
        arrow = slide.shapes.add_connector(1, Inches(6.5), Inches(2.7), Inches(6.5), Inches(4.2))
        arrow.line.color.rgb = RGBColor(30, 64, 124)
        arrow.line.width = Pt(4)
        
        # Satellite
        sat_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(2), Inches(4.5),
                                        Inches(9), Inches(1.2))
        sat_box.fill.solid()
        sat_box.fill.fore_color.rgb = RGBColor(15, 32, 66)
        sat_box.line.color.rgb = RGBColor(30, 64, 124)
        sat_box.line.width = Pt(3)
        
        sat_text = slide.shapes.add_textbox(Inches(2.2), Inches(4.7), Inches(8.6), Inches(0.8))
        sat_frame = sat_text.text_frame
        sat_frame.text = "SATELLITE - À Bord | OS Thales Alenia | Agent IA"
        sat_frame.paragraphs[0].font.size = Pt(18)
        sat_frame.paragraphs[0].font.bold = True
        sat_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        sat_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # ============================================
    # SLIDE 5: SCENARIO 1
    # ============================================
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        if slide.shapes.title:
            slide.shapes.title.text = "Scénario 1: Détection Autonome d'Anomalies"
    except:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "Scénario 1: Détection Autonome d'Anomalies"
        title_frame.paragraphs[0].font.size = Pt(36)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Workflow
    workflow_image = os.path.join(images_dir, "anomaly_workflow.png")
    if os.path.exists(workflow_image):
        slide.shapes.add_picture(workflow_image, Inches(1), Inches(1.5),
                                width=Inches(11), height=Inches(4.5))
    else:
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.4), Inches(4.5))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        content_frame.text = "Processus Autonome:"
        content_frame.paragraphs[0].font.size = Pt(22)
        content_frame.paragraphs[0].font.bold = True
        content_frame.paragraphs[0].space_after = Pt(12)
        
        steps = [
            "1. Capteur → Données télémétrie en temps réel",
            "2. Agent IA (LLM) → Analyse et détection d'anomalies",
            "3. Moteur Décision → Recommandations autonomes",
            "4. Exécuteur Action → Correction automatique",
            "5. Alerte Sol → Notification si critique"
        ]
        
        for step in steps:
            p = content_frame.add_paragraph()
            p.text = step
            p.font.size = Pt(18)
            p.space_after = Pt(8)
    
    # ============================================
    # SLIDE 6: SCENARIO 2 (Before/After)
    # ============================================
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        if slide.shapes.title:
            slide.shapes.title.text = "Scénario 2: Détection Agricole & Stress Hydrique"
    except:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "Scénario 2: Détection Agricole & Stress Hydrique"
        title_frame.paragraphs[0].font.size = Pt(36)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Before/After images
    before_img = os.path.join(images_dir, "before_processing.jpg")
    after_img = os.path.join(images_dir, "after_processing.jpg")
    
    if os.path.exists(before_img) and os.path.exists(after_img):
        # Side by side
        slide.shapes.add_picture(before_img, Inches(0.8), Inches(1.5),
                                width=Inches(5.5), height=Inches(4))
        slide.shapes.add_picture(after_img, Inches(7.0), Inches(1.5),
                                width=Inches(5.5), height=Inches(4))
        
        # Labels
        before_label = slide.shapes.add_textbox(Inches(0.8), Inches(5.6), Inches(5.5), Inches(0.3))
        before_frame = before_label.text_frame
        before_frame.text = "AVANT: 500MB-2GB"
        before_frame.paragraphs[0].font.size = Pt(16)
        before_frame.paragraphs[0].font.bold = True
        before_frame.paragraphs[0].font.color.rgb = RGBColor(200, 0, 0)
        before_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        after_label = slide.shapes.add_textbox(Inches(7.0), Inches(5.6), Inches(5.5), Inches(0.3))
        after_frame = after_label.text_frame
        after_frame.text = "APRÈS: 1-5MB (99.8% réduction)"
        after_frame.paragraphs[0].font.size = Pt(16)
        after_frame.paragraphs[0].font.bold = True
        after_frame.paragraphs[0].font.color.rgb = RGBColor(0, 150, 0)
        after_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    else:
        # Metrics fallback
        metrics = [
            ("99.8%", "Réduction", "50GB → 100MB/jour"),
            ("150K-1.8M$", "Économies/an", "Par satellite"),
            ("Même Jour", "Alertes", "vs lendemain")
        ]
        
        for i, (big, label, detail) in enumerate(metrics):
            x = 0.8 + i * 4.2
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(x), Inches(1.5),
                                         Inches(3.8), Inches(2))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(70, 130, 180)
            card.line.color.rgb = RGBColor(15, 32, 66)
            card.line.width = Pt(3)
            
            num = slide.shapes.add_textbox(Inches(x + 0.2), Inches(1.7), Inches(3.4), Inches(0.8))
            num_frame = num.text_frame
            num_frame.text = big
            num_frame.paragraphs[0].font.size = Pt(32)
            num_frame.paragraphs[0].font.bold = True
            num_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # ============================================
    # SLIDE 7: CONTACT
    # ============================================
    try:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        if slide.shapes.title:
            slide.shapes.title.text = "Contact"
    except:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(12), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "Contact"
        title_frame.paragraphs[0].font.size = Pt(48)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Logo
    logo_img = os.path.join(images_dir, "logo.png")
    if os.path.exists(logo_img):
        slide.shapes.add_picture(logo_img, Inches(10), Inches(0.3),
                                width=Inches(2.5), height=Inches(0.8))
    
    # Contact info
    contact_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9.4), Inches(3.5))
    contact_frame = contact_box.text_frame
    contact_frame.word_wrap = True
    contact_frame.text = "Alexandre GON"
    contact_frame.paragraphs[0].font.size = Pt(32)
    contact_frame.paragraphs[0].font.bold = True
    contact_frame.paragraphs[0].space_after = Pt(16)
    
    p = contact_frame.add_paragraph()
    p.text = "Architecture & Solutions IA"
    p.font.size = Pt(20)
    p.space_after = Pt(20)
    
    p = contact_frame.add_paragraph()
    p.text = "alex.gon26@gmail.com"
    p.font.size = Pt(18)
    p.space_after = Pt(8)
    
    p = contact_frame.add_paragraph()
    p.text = "alex.gon@eliago.com"
    p.font.size = Pt(18)
    p.space_after = Pt(8)
    
    p = contact_frame.add_paragraph()
    p.text = "+33 6 27 79 37 72"
    p.font.size = Pt(18)
    p.font.bold = True
    
    # Photo
    photo_img = os.path.join(images_dir, "photo.jpg")
    if os.path.exists(photo_img):
        slide.shapes.add_picture(photo_img, Inches(1), Inches(2.5),
                                width=Inches(2), height=Inches(2.5))
    
    return prs

if __name__ == "__main__":
    template_path = "Templates/Modele_Decalage.pptx"
    
    if not os.path.exists(template_path):
        print(f"❌ Template not found: {template_path}")
        print(f"Current directory: {os.getcwd()}")
        exit(1)
    
    print(f"🎨 Génération présentation avec template personnalisé...")
    print(f"📄 Template: {template_path}")
    
    prs = create_presentation_with_template(template_path)
    
    if prs:
        output = "presentation_thales_custom_template.pptx"
        prs.save(output)
        print(f"✓ Présentation créée: {output}")
        print(f"✓ Nombre de slides: {len(prs.slides)}")
        print(f"\n💡 Images utilisées depuis 'images/':")
        images = ["background_space.jpg", "before_processing.jpg", "after_processing.jpg"]
        for img in images:
            if os.path.exists(f"images/{img}"):
                print(f"   ✓ {img}")
            else:
                print(f"   ⚠ {img} (non trouvé)")
    else:
        print("❌ Échec de génération")

