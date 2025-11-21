#!/usr/bin/env python3
"""
Generate Thales presentation CORRECTLY using template placeholders
Respects margins and uses proper slide layouts
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

def add_image_safe(slide, image_path, left, top, width=None, height=None):
    """Add image if exists, respecting margins"""
    if os.path.exists(image_path):
        if width and height:
            return slide.shapes.add_picture(image_path, left, top, width, height)
        else:
            return slide.shapes.add_picture(image_path, left, top)
    return None

def create_presentation_correct(template_path):
    """Create presentation using template placeholders correctly"""
    
    if not os.path.exists(template_path):
        print(f"❌ Template not found: {template_path}")
        return None
    
    print(f"📄 Loading template: {template_path}")
    prs = Presentation(template_path)
    
    images_dir = "images"
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    # ============================================
    # SLIDE 1: TITLE (Layout 0: TITLE)
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # TITLE layout
    
    # Use CENTER_TITLE placeholder (index 0)
    if slide.shapes.title:
        slide.shapes.title.text = "IA Embarquée pour\nSystèmes Satellitaires"
    
    # Use SUBTITLE placeholder (index 1)
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "Une Approche Architecturale pour l'Intelligence Autonome à Bord\n\nAlexandre GON\nArchitecture & Solutions IA"
    
    # Add background image if available (behind placeholders)
    bg_image = os.path.join(images_dir, "background_space.jpg")
    if os.path.exists(bg_image):
        bg = slide.shapes.add_picture(bg_image, Inches(0), Inches(0),
                                     slide_width, slide_height)
        # Move to back (z-order)
        bg.element.getparent().remove(bg.element)
        slide.shapes._spTree.insert(2, bg.element)
    
    # ============================================
    # SLIDE 2: THE CHALLENGE (Layout 2: TITLE_AND_BODY)
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[2])  # TITLE_AND_BODY
    
    # Use TITLE placeholder
    if slide.shapes.title:
        slide.shapes.title.text = "Le Défi"
    
    # Use BODY placeholder (index 1)
    if len(slide.placeholders) > 1:
        body = slide.placeholders[1]
        body.text = "50-500GB/jour: Transmission coûteuse"
        
        # Add more items
        tf = body.text_frame
        p = tf.add_paragraph()
        p.text = "Délais critiques: Alertes retardées"
        p.level = 0
        p.space_after = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "Pannes communication: Pas de décision autonome"
        p.level = 0
        p.space_after = Pt(12)
        
        p = tf.add_paragraph()
        p.text = "150K-1.8M$/an: Coûts opérationnels"
        p.level = 0
        p.font.bold = True
    
    # ============================================
    # SLIDE 3: SOLUTION (Layout 3: TITLE_AND_TWO_COLUMNS)
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[3])  # TITLE_AND_TWO_COLUMNS
    
    if slide.shapes.title:
        slide.shapes.title.text = "La Solution"
    
    # Left column (index 1)
    if len(slide.placeholders) > 1:
        left_col = slide.placeholders[1]
        left_col.text = "Opérations Autonomes - 80%"
        tf = left_col.text_frame
        p = tf.add_paragraph()
        p.text = "Décisions sans intervention sol"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "Réduction Coûts - 150K-1.8M$/an"
        p.level = 0
        p.space_after = Pt(12)
        p = tf.add_paragraph()
        p.text = "Transmission -80-99%"
        p.level = 1
    
    # Right column (index 2)
    if len(slide.placeholders) > 2:
        right_col = slide.placeholders[2]
        right_col.text = "Maintenance Prédictive - Lifetime +"
        tf = right_col.text_frame
        p = tf.add_paragraph()
        p.text = "Détection précoce d'anomalies"
        p.level = 1
        
        p = tf.add_paragraph()
        p.text = "Solution Universelle - Portfolio-wide"
        p.level = 0
        p.space_after = Pt(12)
        p = tf.add_paragraph()
        p.text = "Tous types satellites"
        p.level = 1
    
    # ============================================
    # SLIDE 4: ARCHITECTURE (Layout 4: TITLE_ONLY + image)
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[4])  # TITLE_ONLY
    
    if slide.shapes.title:
        slide.shapes.title.text = "Architecture Système"
    
    # Add diagram image in content area (respecting margins)
    # Margins from template: left 0.90", top 2.18", width 8.21"
    arch_image = os.path.join(images_dir, "architecture_diagram.png")
    if os.path.exists(arch_image):
        # Use placeholder area for image positioning
        # Title is at 0.90" left, 0.92" top, 8.21" wide
        # Image should be below title, same margins
        slide.shapes.add_picture(arch_image, 
                                Inches(0.90), Inches(2.18),  # Same left as title, below it
                                width=Inches(8.21), height=Inches(3.0))
    else:
        # Fallback: Create simple diagram using template margins
        # Earth box (using template margins)
        earth_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(0.90), Inches(2.18),
                                        Inches(8.21), Inches(1.0))
        earth_box.fill.solid()
        earth_box.fill.fore_color.rgb = RGBColor(70, 130, 180)
        earth_box.line.color.rgb = RGBColor(15, 32, 66)
        earth_box.line.width = Pt(2)
        
        earth_text = slide.shapes.add_textbox(Inches(0.90), Inches(2.3), Inches(8.21), Inches(0.8))
        earth_frame = earth_text.text_frame
        earth_frame.text = "TERRE - Station Terrestre | Backend Central"
        earth_frame.paragraphs[0].font.size = Pt(16)
        earth_frame.paragraphs[0].font.bold = True
        earth_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Arrow
        arrow = slide.shapes.add_connector(1, 
                                         Inches(5.0), Inches(3.2),
                                         Inches(5.0), Inches(4.0))
        arrow.line.color.rgb = RGBColor(30, 64, 124)
        arrow.line.width = Pt(3)
        
        # Satellite
        sat_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(0.90), Inches(4.2),
                                        Inches(8.21), Inches(1.0))
        sat_box.fill.solid()
        sat_box.fill.fore_color.rgb = RGBColor(15, 32, 66)
        sat_box.line.color.rgb = RGBColor(30, 64, 124)
        sat_box.line.width = Pt(2)
        
        sat_text = slide.shapes.add_textbox(Inches(0.90), Inches(4.3), Inches(8.21), Inches(0.8))
        sat_frame = sat_text.text_frame
        sat_frame.text = "SATELLITE - À Bord | OS Thales Alenia | Agent IA"
        sat_frame.paragraphs[0].font.size = Pt(16)
        sat_frame.paragraphs[0].font.bold = True
        sat_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        sat_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # ============================================
    # SLIDE 5: SCENARIO 1 (Layout 2: TITLE_AND_BODY)
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[2])  # TITLE_AND_BODY
    
    if slide.shapes.title:
        slide.shapes.title.text = "Scénario 1: Détection Autonome d'Anomalies"
    
    if len(slide.placeholders) > 1:
        body = slide.placeholders[1]
        body.text = "Processus Autonome:"
        tf = body.text_frame
        tf.paragraphs[0].font.bold = True
        
        steps = [
            "1. Capteur → Données télémétrie en temps réel",
            "2. Agent IA (LLM) → Analyse et détection d'anomalies",
            "3. Moteur Décision → Recommandations autonomes",
            "4. Exécuteur Action → Correction automatique",
            "5. Alerte Sol → Notification si critique"
        ]
        
        for step in steps:
            p = tf.add_paragraph()
            p.text = step
            p.level = 0
            p.space_after = Pt(8)
    
    # Try to add workflow image (if exists, add below content area)
    workflow_image = os.path.join(images_dir, "anomaly_workflow.png")
    if os.path.exists(workflow_image):
        # Add below body placeholder (body ends around 4.86")
        slide.shapes.add_picture(workflow_image, 
                                Inches(0.90), Inches(4.0),
                                width=Inches(8.21), height=Inches(1.5))
    
    # ============================================
    # SLIDE 6: SCENARIO 2 (Layout 3: TWO_COLUMNS for before/after)
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[3])  # TITLE_AND_TWO_COLUMNS
    
    if slide.shapes.title:
        slide.shapes.title.text = "Scénario 2: Détection Agricole & Stress Hydrique"
    
    # Before/After images in columns
    before_img = os.path.join(images_dir, "before_processing.jpg")
    after_img = os.path.join(images_dir, "after_processing.jpg")
    
    if os.path.exists(before_img) and os.path.exists(after_img):
        # Left column: Before image
        # Placeholder 1 position: (0.90", 2.18"), Size: 4.03" x 2.68"
        slide.shapes.add_picture(before_img,
                                Inches(0.90), Inches(2.18),
                                width=Inches(4.03), height=Inches(2.68))
        
        # Right column: After image
        # Placeholder 2 position: (5.07", 2.18"), Size: 4.03" x 2.68"
        slide.shapes.add_picture(after_img,
                                Inches(5.07), Inches(2.18),
                                width=Inches(4.03), height=Inches(2.68))
        
        # Labels using placeholders
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = "AVANT: 500MB-2GB"
        if len(slide.placeholders) > 2:
            slide.placeholders[2].text = "APRÈS: 1-5MB\n(99.8% réduction)"
    else:
        # Fallback: Use placeholders for text
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = "99.8% Réduction\n50GB → 100MB/jour"
        if len(slide.placeholders) > 2:
            slide.placeholders[2].text = "150K-1.8M$/an\nÉconomies par satellite"
    
    # ============================================
    # SLIDE 7: CONTACT (Layout 2: TITLE_AND_BODY)
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[2])  # TITLE_AND_BODY
    
    if slide.shapes.title:
        slide.shapes.title.text = "Contact"
    
    if len(slide.placeholders) > 1:
        body = slide.placeholders[1]
        body.text = "Alexandre GON"
        tf = body.text_frame
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.size = Pt(24)
        
        p = tf.add_paragraph()
        p.text = "Architecture & Solutions IA"
        p.space_after = Pt(16)
        
        p = tf.add_paragraph()
        p.text = "alex.gon26@gmail.com"
        p.space_after = Pt(8)
        
        p = tf.add_paragraph()
        p.text = "alex.gon@eliago.com"
        p.space_after = Pt(8)
        
        p = tf.add_paragraph()
        p.text = "+33 6 27 79 37 72"
        p.font.bold = True
    
    # Add logo/photo if available (in margin area)
    logo_img = os.path.join(images_dir, "logo.png")
    if os.path.exists(logo_img):
        # Add in top-right corner (respecting margins)
        slide.shapes.add_picture(logo_img, Inches(8.5), Inches(0.3),
                                width=Inches(1.2), height=Inches(0.5))
    
    photo_img = os.path.join(images_dir, "photo.jpg")
    if os.path.exists(photo_img):
        # Add on left side, below title area
        slide.shapes.add_picture(photo_img, Inches(0.90), Inches(2.5),
                                width=Inches(2), height=Inches(2))
    
    return prs

if __name__ == "__main__":
    template_path = "Templates/Modele_Decalage.pptx"
    
    if not os.path.exists(template_path):
        print(f"❌ Template not found: {template_path}")
        exit(1)
    
    print(f"🎨 Génération présentation avec template (méthode correcte)...")
    print(f"📄 Template: {template_path}")
    
    prs = create_presentation_correct(template_path)
    
    if prs:
        output = "presentation_thales_template_correct.pptx"
        prs.save(output)
        print(f"✓ Présentation créée: {output}")
        print(f"✓ Nombre de slides: {len(prs.slides)}")
        print(f"\n✅ Utilise les placeholders du template")
        print(f"✅ Respecte les marges définies")
        print(f"✅ Utilise les layouts appropriés")
    else:
        print("❌ Échec de génération")

