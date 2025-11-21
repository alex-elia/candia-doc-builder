#!/usr/bin/env python3
"""
Generate enhanced PPTX presentation with visual templates and images
for Thales Alenia Space pitch
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR
import os

# Premium color scheme
NAVY = RGBColor(15, 32, 66)
BLUE = RGBColor(30, 64, 124)
LIGHT_BLUE = RGBColor(70, 130, 180)
GOLD = RGBColor(184, 134, 11)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
GRAY = RGBColor(128, 128, 128)

def add_image_safe(slide, image_path, left, top, width=None, height=None):
    """Add image to slide if file exists, otherwise return None"""
    if os.path.exists(image_path):
        if width and height:
            return slide.shapes.add_picture(image_path, left, top, width, height)
        else:
            return slide.shapes.add_picture(image_path, left, top)
    else:
        print(f"⚠️  Image not found: {image_path}")
        return None

def create_enhanced_presentation(template_path=None, images_dir="images"):
    """
    Create enhanced presentation with visual templates and images
    
    Args:
        template_path: Path to existing PowerPoint template (.potx or .pptx)
        images_dir: Directory containing images to use
    """
    
    # Load template if provided, otherwise create new
    if template_path and os.path.exists(template_path):
        print(f"📄 Loading template: {template_path}")
        prs = Presentation(template_path)
    else:
        prs = Presentation()
        # Set 16:9 aspect ratio
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    
    # Create images directory if it doesn't exist
    if not os.path.exists(images_dir):
        os.makedirs(images_dir)
        print(f"📁 Created images directory: {images_dir}")
    
    # ============================================
    # SLIDE 1: TITLE SLIDE (with background image)
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    # Try to add background image
    bg_image = os.path.join(images_dir, "background_space.jpg")
    if os.path.exists(bg_image):
        # Add as background (full slide)
        slide.shapes.add_picture(bg_image, Inches(0), Inches(0), 
                                prs.slide_width, prs.slide_height)
        # Add semi-transparent overlay for readability
        overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                                        Inches(0), Inches(0),
                                        prs.slide_width, prs.slide_height)
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
        overlay.fill.transparency = 0.6
        overlay.line.fill.background()
    
    title.text = "IA Embarquée pour\nSystèmes Satellitaires"
    title_frame = title.text_frame
    title_frame.paragraphs[0].font.size = Pt(48)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WHITE if os.path.exists(bg_image) else NAVY
    
    subtitle.text = "Une Approche Architecturale pour l'Intelligence Autonome à Bord\n\nAlexandre GON\nArchitecture & Solutions IA"
    subtitle_frame = subtitle.text_frame
    subtitle_frame.paragraphs[0].font.size = Pt(20)
    subtitle_frame.paragraphs[0].font.color.rgb = WHITE if os.path.exists(bg_image) else BLUE
    
    # ============================================
    # SLIDE 2: THE CHALLENGE (with icons/images)
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Le Défi"
    title_frame.paragraphs[0].font.size = Pt(52)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = NAVY
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Challenge items with icons
    challenges = [
        ("50-500GB/jour", "Transmission coûteuse", "icon_cost.png"),
        ("Délais critiques", "Alertes retardées", "icon_time.png"),
        ("Pannes communication", "Pas de décision autonome", "icon_network.png"),
        ("150K-1.8M$/an", "Coûts opérationnels", "icon_money.png")
    ]
    
    y_start = 1.8
    for i, (metric, desc, icon) in enumerate(challenges):
        y_pos = y_start + i * 1.1
        x_icon = 1.0
        x_text = 2.5
        
        # Icon (if available)
        icon_path = os.path.join(images_dir, icon)
        if os.path.exists(icon_path):
            slide.shapes.add_picture(icon_path, Inches(x_icon), Inches(y_pos),
                                    width=Inches(1), height=Inches(0.8))
        
        # Metric (bold, large)
        metric_box = slide.shapes.add_textbox(Inches(x_text), Inches(y_pos), 
                                             Inches(4), Inches(0.4))
        metric_frame = metric_box.text_frame
        metric_frame.text = metric
        metric_frame.paragraphs[0].font.size = Pt(24)
        metric_frame.paragraphs[0].font.bold = True
        metric_frame.paragraphs[0].font.color.rgb = NAVY
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(x_text), Inches(y_pos + 0.5), 
                                           Inches(8), Inches(0.3))
        desc_frame = desc_box.text_frame
        desc_frame.text = desc
        desc_frame.paragraphs[0].font.size = Pt(16)
        desc_frame.paragraphs[0].font.color.rgb = GRAY
    
    # ============================================
    # SLIDE 3: SOLUTION (with visual boxes)
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "La Solution"
    title_frame.paragraphs[0].font.size = Pt(52)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = NAVY
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Value props in visual cards
    value_props = [
        ("Opérations\nAutonomes", "80%", "Décisions sans intervention sol", LIGHT_BLUE),
        ("Réduction\nCoûts", "150K-1.8M$/an", "Transmission -80-99%", BLUE),
        ("Maintenance\nPrédictive", "Lifetime +", "Détection précoce", NAVY),
        ("Solution\nUniverselle", "Portfolio-wide", "Tous types satellites", GOLD)
    ]
    
    # 2x2 grid
    positions = [
        (0.8, 1.8),   # Top left
        (7.0, 1.8),   # Top right
        (0.8, 4.0),   # Bottom left
        (7.0, 4.0)    # Bottom right
    ]
    
    for i, ((title, metric, desc, color), (x, y)) in enumerate(zip(value_props, positions)):
        # Card background
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(x), Inches(y),
                                     Inches(5.5), Inches(2))
        card.fill.solid()
        card.fill.fore_color.rgb = color
        card.line.color.rgb = BLUE
        card.line.width = Pt(3)
        
        # Title
        title_text = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.2),
                                            Inches(5.1), Inches(0.6))
        title_frame = title_text.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(20)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = WHITE if color == NAVY else NAVY
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Metric (big, gold)
        metric_text = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.9),
                                              Inches(5.1), Inches(0.6))
        metric_frame = metric_text.text_frame
        metric_frame.text = metric
        metric_frame.paragraphs[0].font.size = Pt(32)
        metric_frame.paragraphs[0].font.bold = True
        metric_frame.paragraphs[0].font.color.rgb = GOLD
        metric_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Description
        desc_text = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 1.6),
                                            Inches(5.1), Inches(0.3))
        desc_frame = desc_text.text_frame
        desc_frame.text = desc
        desc_frame.paragraphs[0].font.size = Pt(14)
        desc_frame.paragraphs[0].font.color.rgb = WHITE if color == NAVY else BLACK
        desc_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # ============================================
    # SLIDE 4: ARCHITECTURE (with diagram image)
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Architecture Système"
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = NAVY
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Try to add architecture diagram
    arch_image = os.path.join(images_dir, "architecture_diagram.png")
    if os.path.exists(arch_image):
        # Add diagram (centered, large)
        slide.shapes.add_picture(arch_image, Inches(1.5), Inches(1.5),
                                width=Inches(10), height=Inches(5))
    else:
        # Fallback: Create simple diagram with shapes
        # Earth box
        earth_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(2), Inches(1.5),
                                        Inches(9), Inches(1.2))
        earth_box.fill.solid()
        earth_box.fill.fore_color.rgb = LIGHT_BLUE
        earth_box.line.color.rgb = BLUE
        earth_box.line.width = Pt(3)
        
        earth_text = slide.shapes.add_textbox(Inches(2.2), Inches(1.7), Inches(8.6), Inches(0.8))
        earth_frame = earth_text.text_frame
        earth_frame.text = "TERRE - Station Terrestre"
        earth_frame.paragraphs[0].font.size = Pt(18)
        earth_frame.paragraphs[0].font.bold = True
        earth_frame.paragraphs[0].font.color.rgb = NAVY
        earth_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Arrow
        arrow = slide.shapes.add_connector(1, Inches(6.5), Inches(2.7), Inches(6.5), Inches(4.2))
        arrow.line.color.rgb = BLUE
        arrow.line.width = Pt(4)
        
        # Satellite box
        sat_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(2), Inches(4.5),
                                        Inches(9), Inches(1.2))
        sat_box.fill.solid()
        sat_box.fill.fore_color.rgb = NAVY
        sat_box.line.color.rgb = BLUE
        sat_box.line.width = Pt(3)
        
        sat_text = slide.shapes.add_textbox(Inches(2.2), Inches(4.7), Inches(8.6), Inches(0.8))
        sat_frame = sat_text.text_frame
        sat_frame.text = "SATELLITE - À Bord"
        sat_frame.paragraphs[0].font.size = Pt(18)
        sat_frame.paragraphs[0].font.bold = True
        sat_frame.paragraphs[0].font.color.rgb = WHITE
        sat_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # ============================================
    # SLIDE 5: SCENARIO 1 (with workflow image)
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Scénario 1: Détection Autonome d'Anomalies"
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = NAVY
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Try to add workflow diagram
    workflow_image = os.path.join(images_dir, "anomaly_workflow.png")
    if os.path.exists(workflow_image):
        slide.shapes.add_picture(workflow_image, Inches(1), Inches(1.5),
                                width=Inches(11), height=Inches(4.5))
    else:
        # Fallback: Text description
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.4), Inches(4.5))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        content_frame.text = "Processus:"
        content_frame.paragraphs[0].font.size = Pt(20)
        content_frame.paragraphs[0].font.bold = True
        content_frame.paragraphs[0].font.color.rgb = NAVY
        content_frame.paragraphs[0].space_after = Pt(12)
        
        steps = [
            "1. Capteur → Données télémétrie",
            "2. Agent IA (LLM) → Analyse en temps réel",
            "3. Moteur Décision → Recommandations",
            "4. Exécuteur Action → Correction autonome",
            "5. Alerte Sol → Si critique"
        ]
        
        for step in steps:
            p = content_frame.add_paragraph()
            p.text = step
            p.font.size = Pt(16)
            p.space_after = Pt(8)
    
    # ============================================
    # SLIDE 6: SCENARIO 2 (with before/after images)
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Scénario 2: Détection Agricole & Stress Hydrique"
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = NAVY
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Before/After comparison
    before_image = os.path.join(images_dir, "before_processing.jpg")
    after_image = os.path.join(images_dir, "after_processing.jpg")
    
    if os.path.exists(before_image) and os.path.exists(after_image):
        # Side by side comparison
        slide.shapes.add_picture(before_image, Inches(0.8), Inches(1.5),
                                width=Inches(5.5), height=Inches(4))
        slide.shapes.add_picture(after_image, Inches(7.0), Inches(1.5),
                                width=Inches(5.5), height=Inches(4))
        
        # Labels
        before_label = slide.shapes.add_textbox(Inches(0.8), Inches(5.6), Inches(5.5), Inches(0.3))
        before_frame = before_label.text_frame
        before_frame.text = "AVANT: 500MB-2GB par image"
        before_frame.paragraphs[0].font.size = Pt(14)
        before_frame.paragraphs[0].font.bold = True
        before_frame.paragraphs[0].font.color.rgb = RGBColor(200, 0, 0)
        before_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        after_label = slide.shapes.add_textbox(Inches(7.0), Inches(5.6), Inches(5.5), Inches(0.3))
        after_frame = after_label.text_frame
        after_frame.text = "APRÈS: 1-5MB rapport"
        after_frame.paragraphs[0].font.size = Pt(14)
        after_frame.paragraphs[0].font.bold = True
        after_frame.paragraphs[0].font.color.rgb = RGBColor(0, 150, 0)
        after_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    else:
        # Fallback: Metrics
        metrics = [
            ("99.8%", "Réduction Transmission", "50GB/jour → 100MB/jour"),
            ("150K-1.8M$", "Économies/An", "Par satellite"),
            ("Même Jour", "Alertes vs", "Traitement lendemain")
        ]
        
        for i, (big_num, label, detail) in enumerate(metrics):
            x_pos = 0.8 + i * 4.2
            
            metric_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                               Inches(x_pos), Inches(1.5),
                                               Inches(3.8), Inches(2))
            metric_box.fill.solid()
            metric_box.fill.fore_color.rgb = LIGHT_BLUE
            metric_box.line.color.rgb = BLUE
            metric_box.line.width = Pt(3)
            
            num_text = slide.shapes.add_textbox(Inches(x_pos + 0.2), Inches(1.7),
                                               Inches(3.4), Inches(0.8))
            num_frame = num_text.text_frame
            num_frame.text = big_num
            num_frame.paragraphs[0].font.size = Pt(32)
            num_frame.paragraphs[0].font.bold = True
            num_frame.paragraphs[0].font.color.rgb = GOLD
            num_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # ============================================
    # SLIDE 7: CONTACT (with logo/photo)
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Try to add logo/photo
    logo_image = os.path.join(images_dir, "logo.png")
    photo_image = os.path.join(images_dir, "photo.jpg")
    
    if os.path.exists(logo_image):
        slide.shapes.add_picture(logo_image, Inches(10), Inches(0.3),
                                width=Inches(2.5), height=Inches(0.8))
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(12), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Contact"
    title_frame.paragraphs[0].font.size = Pt(48)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = NAVY
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Contact info with photo
    if os.path.exists(photo_image):
        slide.shapes.add_picture(photo_image, Inches(1), Inches(2.5),
                                width=Inches(2), height=Inches(2.5))
        contact_x = 3.5
    else:
        contact_x = 2
    
    contact_box = slide.shapes.add_textbox(Inches(contact_x), Inches(2.5),
                                          Inches(9.4), Inches(3.5))
    contact_frame = contact_box.text_frame
    contact_frame.word_wrap = True
    contact_frame.text = "Alexandre GON"
    contact_frame.paragraphs[0].font.size = Pt(32)
    contact_frame.paragraphs[0].font.bold = True
    contact_frame.paragraphs[0].font.color.rgb = NAVY
    contact_frame.paragraphs[0].space_after = Pt(16)
    
    p = contact_frame.add_paragraph()
    p.text = "Architecture & Solutions IA"
    p.font.size = Pt(20)
    p.font.color.rgb = BLUE
    p.space_after = Pt(20)
    
    p = contact_frame.add_paragraph()
    p.text = "alex.gon26@gmail.com"
    p.font.size = Pt(18)
    p.space_after = Pt(8)
    
    p = contact_frame.add_paragraph()
    p.text = "alex.gon@eliago.com"
    p.font.size = Pt(18)
    p.space_after = Pt(8)
    
    p = contact_frame.add_paragraph()
    p.text = "+33 6 27 79 37 72"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = NAVY
    
    return prs

if __name__ == "__main__":
    import sys
    
    # Check for template argument
    template = None
    if len(sys.argv) > 1:
        template = sys.argv[1]
        if not os.path.exists(template):
            print(f"⚠️  Template not found: {template}")
            template = None
    
    print("🎨 Génération de la présentation visuelle Thales Alenia Space...")
    presentation = create_enhanced_presentation(template_path=template)
    output_file = "presentation_thales_alenia_enhanced.pptx"
    presentation.save(output_file)
    print(f"✓ Présentation créée: {output_file}")
    print(f"✓ Nombre de slides: {len(presentation.slides)}")
    print(f"\n💡 Pour utiliser des images, placez-les dans le dossier 'images/':")
    print(f"   - background_space.jpg (slide 1)")
    print(f"   - architecture_diagram.png (slide 4)")
    print(f"   - anomaly_workflow.png (slide 5)")
    print(f"   - before_processing.jpg / after_processing.jpg (slide 6)")
    print(f"   - logo.png (slide 7)")
    print(f"   - photo.jpg (slide 7)")

