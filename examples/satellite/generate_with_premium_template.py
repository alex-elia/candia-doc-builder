#!/usr/bin/env python3
"""
Generate Thales presentation with premium tech template and images
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Premium Tech Sustainable Colors
DARK_GREEN = RGBColor(0, 100, 80)
TECH_BLUE = RGBColor(0, 120, 215)
LIGHT_GREEN = RGBColor(144, 238, 144)
DARK_GRAY = RGBColor(30, 30, 30)
LIGHT_GRAY = RGBColor(240, 240, 240)
WHITE = RGBColor(255, 255, 255)
ACCENT_GOLD = RGBColor(255, 215, 0)
NAVY = RGBColor(15, 32, 66)

def add_image_safe(slide, image_path, left, top, width=None, height=None):
    """Add image if exists"""
    if os.path.exists(image_path):
        if width and height:
            return slide.shapes.add_picture(image_path, left, top, width, height)
        else:
            return slide.shapes.add_picture(image_path, left, top)
    return None

def create_premium_presentation():
    """Create presentation with premium tech template"""
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    images_dir = "images"
    
    # ============================================
    # SLIDE 1: TITLE (Premium with background)
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Dark background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                               Inches(0), Inches(0),
                               prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_GRAY
    bg.line.fill.background()
    
    # Accent line top
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0), Inches(0),
                                   prs.slide_width, Inches(0.2))
    accent.fill.solid()
    accent.fill.fore_color.rgb = TECH_BLUE
    accent.line.fill.background()
    
    # Try background image
    bg_image = os.path.join(images_dir, "background_space.jpg")
    if os.path.exists(bg_image):
        img = slide.shapes.add_picture(bg_image, Inches(0), Inches(0),
                                      prs.slide_width, prs.slide_height)
        # Dark overlay for readability
        overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Inches(0), Inches(0),
                                        prs.slide_width, prs.slide_height)
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
        overlay.fill.transparency = 0.5
        overlay.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.3), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "IA Embarquée pour\nSystèmes Satellitaires"
    title_frame.paragraphs[0].font.size = Pt(52)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WHITE
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.3), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Une Approche Architecturale pour l'Intelligence Autonome à Bord\n\nAlexandre GON • Architecture & Solutions IA"
    subtitle_frame.paragraphs[0].font.size = Pt(20)
    subtitle_frame.paragraphs[0].font.color.rgb = LIGHT_GREEN
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # ============================================
    # SLIDE 2: THE CHALLENGE (Visual cards)
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Light background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                                               prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_GRAY
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Le Défi"
    title_frame.paragraphs[0].font.size = Pt(48)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = DARK_GREEN
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Challenge cards
    challenges = [
        ("50-500GB/jour", "Transmission coûteuse", DARK_GREEN),
        ("Délais critiques", "Alertes retardées", TECH_BLUE),
        ("Pannes communication", "Pas de décision autonome", NAVY),
        ("150K-1.8M$/an", "Coûts opérationnels", RGBColor(200, 0, 0))
    ]
    
    y_start = 1.8
    for i, (metric, desc, color) in enumerate(challenges):
        y_pos = y_start + i * 1.2
        x_pos = 1.0
        
        # Card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(x_pos), Inches(y_pos),
                                     Inches(11), Inches(1))
        card.fill.solid()
        card.fill.fore_color.rgb = color
        card.line.color.rgb = DARK_GREEN
        card.line.width = Pt(2)
        
        # Metric
        metric_box = slide.shapes.add_textbox(Inches(x_pos + 0.3), Inches(y_pos + 0.1),
                                             Inches(4), Inches(0.4))
        metric_frame = metric_box.text_frame
        metric_frame.text = metric
        metric_frame.paragraphs[0].font.size = Pt(24)
        metric_frame.paragraphs[0].font.bold = True
        metric_frame.paragraphs[0].font.color.rgb = WHITE
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(x_pos + 0.3), Inches(y_pos + 0.55),
                                           Inches(10), Inches(0.3))
        desc_frame = desc_box.text_frame
        desc_frame.text = desc
        desc_frame.paragraphs[0].font.size = Pt(16)
        desc_frame.paragraphs[0].font.color.rgb = WHITE
    
    # ============================================
    # SLIDE 3: SOLUTION (Premium cards)
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                               prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "La Solution"
    title_frame.paragraphs[0].font.size = Pt(48)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = DARK_GREEN
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Value props in 2x2 grid
    value_props = [
        ("Opérations\nAutonomes", "80%", LIGHT_GREEN, DARK_GREEN),
        ("Réduction\nCoûts", "150K-1.8M$/an", TECH_BLUE, WHITE),
        ("Maintenance\nPrédictive", "Lifetime +", DARK_GREEN, WHITE),
        ("Solution\nUniverselle", "Portfolio-wide", ACCENT_GOLD, DARK_GRAY)
    ]
    
    positions = [(0.8, 1.8), (7.0, 1.8), (0.8, 4.0), (7.0, 4.0)]
    
    for i, ((title, metric, bg_color, text_color), (x, y)) in enumerate(zip(value_props, positions)):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(x), Inches(y),
                                     Inches(5.5), Inches(2))
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        card.line.color.rgb = DARK_GREEN
        card.line.width = Pt(3)
        card.shadow.inherit = False
        
        # Title
        title_text = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.2),
                                            Inches(5.1), Inches(0.6))
        title_frame = title_text.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(20)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = text_color
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Metric
        metric_text = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.9),
                                              Inches(5.1), Inches(0.6))
        metric_frame = metric_text.text_frame
        metric_frame.text = metric
        metric_frame.paragraphs[0].font.size = Pt(32)
        metric_frame.paragraphs[0].font.bold = True
        metric_frame.paragraphs[0].font.color.rgb = ACCENT_GOLD if bg_color != ACCENT_GOLD else DARK_GRAY
        metric_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # ============================================
    # SLIDE 4: ARCHITECTURE (with diagram)
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                               prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Architecture Système"
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = DARK_GREEN
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Try architecture diagram
    arch_image = os.path.join(images_dir, "architecture_diagram.png")
    if os.path.exists(arch_image):
        slide.shapes.add_picture(arch_image, Inches(1.5), Inches(1.5),
                                width=Inches(10), height=Inches(5))
    else:
        # Fallback diagram
        earth_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(2), Inches(1.5),
                                        Inches(9), Inches(1.2))
        earth_box.fill.solid()
        earth_box.fill.fore_color.rgb = LIGHT_GREEN
        earth_box.line.color.rgb = DARK_GREEN
        earth_box.line.width = Pt(3)
        
        earth_text = slide.shapes.add_textbox(Inches(2.2), Inches(1.7), Inches(8.6), Inches(0.8))
        earth_frame = earth_text.text_frame
        earth_frame.text = "TERRE - Station Terrestre | Backend Central"
        earth_frame.paragraphs[0].font.size = Pt(18)
        earth_frame.paragraphs[0].font.bold = True
        earth_frame.paragraphs[0].font.color.rgb = DARK_GREEN
        earth_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Arrow
        arrow = slide.shapes.add_connector(1, Inches(6.5), Inches(2.7), Inches(6.5), Inches(4.2))
        arrow.line.color.rgb = TECH_BLUE
        arrow.line.width = Pt(4)
        
        # Satellite
        sat_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(2), Inches(4.5),
                                        Inches(9), Inches(1.2))
        sat_box.fill.solid()
        sat_box.fill.fore_color.rgb = TECH_BLUE
        sat_box.line.color.rgb = DARK_GREEN
        sat_box.line.width = Pt(3)
        
        sat_text = slide.shapes.add_textbox(Inches(2.2), Inches(4.7), Inches(8.6), Inches(0.8))
        sat_frame = sat_text.text_frame
        sat_frame.text = "SATELLITE - À Bord | OS Thales Alenia | Agent IA"
        sat_frame.paragraphs[0].font.size = Pt(18)
        sat_frame.paragraphs[0].font.bold = True
        sat_frame.paragraphs[0].font.color.rgb = WHITE
        sat_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # ============================================
    # SLIDE 5: SCENARIO 1
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                               prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT_GRAY
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Scénario 1: Détection Autonome d'Anomalies"
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = DARK_GREEN
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Workflow image or text
    workflow_image = os.path.join(images_dir, "anomaly_workflow.png")
    if os.path.exists(workflow_image):
        slide.shapes.add_picture(workflow_image, Inches(1), Inches(1.5),
                                width=Inches(11), height=Inches(4.5))
    else:
        # Text workflow
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.4), Inches(4.5))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        content_frame.text = "Processus Autonome:"
        content_frame.paragraphs[0].font.size = Pt(22)
        content_frame.paragraphs[0].font.bold = True
        content_frame.paragraphs[0].font.color.rgb = DARK_GREEN
        content_frame.paragraphs[0].space_after = Pt(12)
        
        steps = [
            "1. Capteur → Données télémétrie en temps réel",
            "2. Agent IA (LLM) → Analyse et détection d'anomalies",
            "3. Moteur Décision → Recommandations autonomes",
            "4. Exécuteur Action → Correction automatique",
            "5. Alerte Sol → Notification si critique"
        ]
        
        for step in steps:
            p = content_frame.add_paragraph()
            p.text = step
            p.font.size = Pt(18)
            p.font.color.rgb = NAVY
            p.space_after = Pt(8)
    
    # ============================================
    # SLIDE 6: SCENARIO 2 (Before/After)
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                               prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Scénario 2: Détection Agricole & Stress Hydrique"
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = DARK_GREEN
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Before/After images
    before_img = os.path.join(images_dir, "before_processing.jpg")
    after_img = os.path.join(images_dir, "after_processing.jpg")
    
    if os.path.exists(before_img) and os.path.exists(after_img):
        slide.shapes.add_picture(before_img, Inches(0.8), Inches(1.5),
                                width=Inches(5.5), height=Inches(4))
        slide.shapes.add_picture(after_img, Inches(7.0), Inches(1.5),
                                width=Inches(5.5), height=Inches(4))
        
        # Labels
        before_label = slide.shapes.add_textbox(Inches(0.8), Inches(5.6), Inches(5.5), Inches(0.3))
        before_frame = before_label.text_frame
        before_frame.text = "AVANT: 500MB-2GB"
        before_frame.paragraphs[0].font.size = Pt(16)
        before_frame.paragraphs[0].font.bold = True
        before_frame.paragraphs[0].font.color.rgb = RGBColor(200, 0, 0)
        before_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        after_label = slide.shapes.add_textbox(Inches(7.0), Inches(5.6), Inches(5.5), Inches(0.3))
        after_frame = after_label.text_frame
        after_frame.text = "APRÈS: 1-5MB (99.8% réduction)"
        after_frame.paragraphs[0].font.size = Pt(16)
        after_frame.paragraphs[0].font.bold = True
        after_frame.paragraphs[0].font.color.rgb = DARK_GREEN
        after_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    else:
        # Metrics fallback
        metrics = [
            ("99.8%", "Réduction", "50GB → 100MB/jour"),
            ("150K-1.8M$", "Économies/an", "Par satellite"),
            ("Même Jour", "Alertes", "vs lendemain")
        ]
        
        for i, (big, label, detail) in enumerate(metrics):
            x = 0.8 + i * 4.2
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(x), Inches(1.5),
                                         Inches(3.8), Inches(2))
            card.fill.solid()
            card.fill.fore_color.rgb = LIGHT_GREEN
            card.line.color.rgb = DARK_GREEN
            card.line.width = Pt(3)
            
            num = slide.shapes.add_textbox(Inches(x + 0.2), Inches(1.7), Inches(3.4), Inches(0.8))
            num_frame = num.text_frame
            num_frame.text = big
            num_frame.paragraphs[0].font.size = Pt(32)
            num_frame.paragraphs[0].font.bold = True
            num_frame.paragraphs[0].font.color.rgb = DARK_GREEN
            num_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # ============================================
    # SLIDE 7: CONTACT
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                               prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_GRAY
    bg.line.fill.background()
    
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0), Inches(0),
                                   prs.slide_width, Inches(0.2))
    accent.fill.solid()
    accent.fill.fore_color.rgb = TECH_BLUE
    accent.line.fill.background()
    
    # Logo
    logo_img = os.path.join(images_dir, "logo.png")
    if os.path.exists(logo_img):
        slide.shapes.add_picture(logo_img, Inches(10), Inches(0.3),
                                width=Inches(2.5), height=Inches(0.8))
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(12), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Contact"
    title_frame.paragraphs[0].font.size = Pt(48)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = WHITE
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Photo
    photo_img = os.path.join(images_dir, "photo.jpg")
    contact_x = 2
    if os.path.exists(photo_img):
        slide.shapes.add_picture(photo_img, Inches(1), Inches(2.5),
                                width=Inches(2), height=Inches(2.5))
        contact_x = 3.5
    
    # Contact info
    contact_box = slide.shapes.add_textbox(Inches(contact_x), Inches(2.5),
                                          Inches(9.4), Inches(3.5))
    contact_frame = contact_box.text_frame
    contact_frame.word_wrap = True
    contact_frame.text = "Alexandre GON"
    contact_frame.paragraphs[0].font.size = Pt(32)
    contact_frame.paragraphs[0].font.bold = True
    contact_frame.paragraphs[0].font.color.rgb = LIGHT_GREEN
    contact_frame.paragraphs[0].space_after = Pt(16)
    
    p = contact_frame.add_paragraph()
    p.text = "Architecture & Solutions IA"
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.space_after = Pt(20)
    
    p = contact_frame.add_paragraph()
    p.text = "alex.gon26@gmail.com"
    p.font.size = Pt(18)
    p.font.color.rgb = LIGHT_GREEN
    p.space_after = Pt(8)
    
    p = contact_frame.add_paragraph()
    p.text = "alex.gon@eliago.com"
    p.font.size = Pt(18)
    p.font.color.rgb = LIGHT_GREEN
    p.space_after = Pt(8)
    
    p = contact_frame.add_paragraph()
    p.text = "+33 6 27 79 37 72"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GOLD
    
    return prs

if __name__ == "__main__":
    print("🎨 Génération présentation premium tech durable...")
    prs = create_premium_presentation()
    output = "presentation_thales_premium.pptx"
    prs.save(output)
    print(f"✓ Présentation créée: {output}")
    print(f"✓ Nombre de slides: {len(prs.slides)}")
    print(f"\n💡 Pour ajouter des images, placez-les dans 'images/':")
    print(f"   - background_space.jpg")
    print(f"   - architecture_diagram.png")
    print(f"   - anomaly_workflow.png")
    print(f"   - before_processing.jpg / after_processing.jpg")
    print(f"   - logo.png / photo.jpg")

