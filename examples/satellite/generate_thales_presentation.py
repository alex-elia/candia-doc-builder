#!/usr/bin/env python3
"""
Generate impactful PPTX presentation for Thales Alenia Space pitch
Based on the pitch document content
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_THEME_COLOR

# Premium color scheme (matching LaTeX document)
NAVY = RGBColor(15, 32, 66)      # Deep professional navy
BLUE = RGBColor(30, 64, 124)     # Rich blue
LIGHT_BLUE = RGBColor(70, 130, 180)  # Light blue for highlights
GOLD = RGBColor(184, 134, 11)    # Elegant gold accent
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
GRAY = RGBColor(128, 128, 128)

def create_presentation():
    """Create the Thales Alenia Space pitch presentation"""
    
    # Create presentation with 16:9 aspect ratio
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Set slide size to 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # ============================================
    # SLIDE 1: TITLE SLIDE
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "IA Embarquée pour\nSystèmes Satellitaires"
    title_frame = title.text_frame
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = NAVY
    
    subtitle.text = "Une Approche Architecturale pour l'Intelligence Autonome à Bord\n\nAlexandre GON\nArchitecture & Solutions IA"
    subtitle_frame = subtitle.text_frame
    subtitle_frame.paragraphs[0].font.size = Pt(18)
    subtitle_frame.paragraphs[0].font.color.rgb = BLUE
    
    # ============================================
    # SLIDE 2: HOOK - The Challenge
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Le Défi"
    title_frame.paragraphs[0].font.size = Pt(48)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = NAVY
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(4.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    p = content_frame.paragraphs[0]
    p.text = "Les satellites modernes génèrent des volumes massifs de données"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = BLUE
    p.space_after = Pt(12)
    
    p = content_frame.add_paragraph()
    p.text = "• Transmission coûteuse: 50-500GB/jour par satellite"
    p.font.size = Pt(20)
    p.font.color.rgb = BLACK
    p.space_after = Pt(8)
    p.level = 0
    
    p = content_frame.add_paragraph()
    p.text = "• Délais critiques: Alertes agricoles/militaires retardées"
    p.font.size = Pt(20)
    p.font.color.rgb = BLACK
    p.space_after = Pt(8)
    p.level = 0
    
    p = content_frame.add_paragraph()
    p.text = "• Pannes de communication: Pas de décision autonome"
    p.font.size = Pt(20)
    p.font.color.rgb = BLACK
    p.space_after = Pt(8)
    p.level = 0
    
    p = content_frame.add_paragraph()
    p.text = "• Coûts opérationnels: 150K-1.8M$/an par satellite"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = NAVY
    p.space_after = Pt(12)
    p.level = 0
    
    # ============================================
    # SLIDE 3: THE SOLUTION - Value Proposition
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "La Solution"
    title_frame.paragraphs[0].font.size = Pt(48)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = NAVY
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Value propositions in boxes
    y_start = 1.8
    box_height = 0.9
    spacing = 0.2
    
    value_props = [
        ("Opérations Autonomes", "Décisions sans intervention sol", "80%"),
        ("Réduction Coûts", "Transmission réduite de 80-99%", "150K-1.8M$/an"),
        ("Maintenance Prédictive", "Détection précoce d'anomalies", "Lifetime +"),
        ("Solution Universelle", "Observation, Communications, Défense", "Portfolio-wide")
    ]
    
    for i, (title, desc, metric) in enumerate(value_props):
        y_pos = y_start + i * (box_height + spacing)
        
        # Box background
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, 
                                    Inches(0.8), Inches(y_pos),
                                    Inches(5.5), Inches(box_height))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BLUE
        box.line.color.rgb = BLUE
        box.line.width = Pt(2)
        
        # Title
        text_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos + 0.1), 
                                           Inches(5), Inches(0.35))
        text_frame = text_box.text_frame
        text_frame.text = title
        text_frame.paragraphs[0].font.size = Pt(18)
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].font.color.rgb = NAVY
        
        # Description
        desc_box = slide.shapes.add_textbox(Inches(1), Inches(y_pos + 0.45), 
                                           Inches(5), Inches(0.3))
        desc_frame = desc_box.text_frame
        desc_frame.text = desc
        desc_frame.paragraphs[0].font.size = Pt(14)
        desc_frame.paragraphs[0].font.color.rgb = BLACK
        
        # Metric (right side)
        metric_box = slide.shapes.add_textbox(Inches(6.5), Inches(y_pos + 0.25), 
                                             Inches(2.5), Inches(0.4))
        metric_frame = metric_box.text_frame
        metric_frame.text = metric
        metric_frame.paragraphs[0].font.size = Pt(20)
        metric_frame.paragraphs[0].font.bold = True
        metric_frame.paragraphs[0].font.color.rgb = GOLD
        metric_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # ============================================
    # SLIDE 4: Architecture Overview
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Architecture Système"
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = NAVY
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Earth/Ground Station box
    earth_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(2), Inches(1.5),
                                      Inches(9), Inches(1.2))
    earth_box.fill.solid()
    earth_box.fill.fore_color.rgb = LIGHT_BLUE
    earth_box.line.color.rgb = BLUE
    earth_box.line.width = Pt(3)
    
    earth_text = slide.shapes.add_textbox(Inches(2.2), Inches(1.7), Inches(8.6), Inches(0.8))
    earth_frame = earth_text.text_frame
    earth_frame.text = "TERRE - Station Terrestre\nBackend Central | Mission Control | Data Processing"
    earth_frame.paragraphs[0].font.size = Pt(18)
    earth_frame.paragraphs[0].font.bold = True
    earth_frame.paragraphs[0].font.color.rgb = NAVY
    earth_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Arrow
    arrow = slide.shapes.add_connector(1, Inches(6.5), Inches(2.7), Inches(6.5), Inches(4.2))
    arrow.line.color.rgb = BLUE
    arrow.line.width = Pt(4)
    
    # Satellite box
    sat_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    Inches(2), Inches(4.5),
                                    Inches(9), Inches(1.2))
    sat_box.fill.solid()
    sat_box.fill.fore_color.rgb = NAVY
    sat_box.line.color.rgb = BLUE
    sat_box.line.width = Pt(3)
    
    sat_text = slide.shapes.add_textbox(Inches(2.2), Inches(4.7), Inches(8.6), Inches(0.8))
    sat_frame = sat_text.text_frame
    sat_frame.text = "SATELLITE - À Bord\nOS Thales Alenia | Agent IA | Capteurs | Communication"
    sat_frame.paragraphs[0].font.size = Pt(18)
    sat_frame.paragraphs[0].font.bold = True
    sat_frame.paragraphs[0].font.color.rgb = WHITE
    sat_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Communication label
    comm_label = slide.shapes.add_textbox(Inches(7.5), Inches(3.2), Inches(2), Inches(0.5))
    comm_frame = comm_label.text_frame
    comm_frame.text = "Liaison\nCommunication"
    comm_frame.paragraphs[0].font.size = Pt(12)
    comm_frame.paragraphs[0].font.color.rgb = BLUE
    comm_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # ============================================
    # SLIDE 5: On-Board AI Architecture
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Architecture IA à Bord"
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = NAVY
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Layers
    layers = [
        ("Embedded Linux / VxWorks RTOS", NAVY, WHITE, 1.5),
        ("Tâches Temps Réel (Priorité 1-10)\nContrôle satellite | Capteurs | Communication", BLUE, WHITE, 2.8),
        ("Conteneur OCI: Inférence IA\nTensorFlow Lite / llama.cpp | LLM Quantifié (500MB)", LIGHT_BLUE, BLACK, 4.1),
        ("Accélération Matérielle (si disponible)\nGPU/NPU pour inférence", GRAY, BLACK, 5.4)
    ]
    
    for i, (text, bg_color, text_color, y_pos) in enumerate(layers):
        layer_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          Inches(1.5), Inches(y_pos),
                                          Inches(10), Inches(0.9))
        layer_box.fill.solid()
        layer_box.fill.fore_color.rgb = bg_color
        layer_box.line.color.rgb = BLUE
        layer_box.line.width = Pt(2)
        
        layer_text = slide.shapes.add_textbox(Inches(1.7), Inches(y_pos + 0.15), 
                                             Inches(9.6), Inches(0.6))
        text_frame = layer_text.text_frame
        text_frame.text = text
        text_frame.paragraphs[0].font.size = Pt(14)
        text_frame.paragraphs[0].font.bold = True
        text_frame.paragraphs[0].font.color.rgb = text_color
        text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # ============================================
    # SLIDE 6: Scenario 1 - Anomaly Detection
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Scénario 1: Détection Autonome d'Anomalies"
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = NAVY
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Challenge box
    challenge_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          Inches(0.8), Inches(1.3),
                                          Inches(5.8), Inches(1.5))
    challenge_box.fill.solid()
    challenge_box.fill.fore_color.rgb = RGBColor(255, 200, 200)  # Light red
    challenge_box.line.color.rgb = RGBColor(200, 0, 0)
    challenge_box.line.width = Pt(2)
    
    challenge_text = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(5.4), Inches(1.1))
    challenge_frame = challenge_text.text_frame
    challenge_frame.text = "DÉFI\nDétecter anomalies (température, communication)\nSans contact sol disponible"
    challenge_frame.paragraphs[0].font.size = Pt(16)
    challenge_frame.paragraphs[0].font.bold = True
    challenge_frame.paragraphs[0].font.color.rgb = RGBColor(150, 0, 0)
    challenge_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Solution box
    solution_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          Inches(6.6), Inches(1.3),
                                          Inches(5.8), Inches(1.5))
    solution_box.fill.solid()
    solution_box.fill.fore_color.rgb = RGBColor(200, 255, 200)  # Light green
    solution_box.line.color.rgb = RGBColor(0, 150, 0)
    solution_box.line.width = Pt(2)
    
    solution_text = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.4), Inches(1.1))
    solution_frame = solution_text.text_frame
    solution_frame.text = "SOLUTION\nAgent LLM surveille télémétrie\nDécision autonome + Alertes"
    solution_frame.paragraphs[0].font.size = Pt(16)
    solution_frame.paragraphs[0].font.bold = True
    solution_frame.paragraphs[0].font.color.rgb = RGBColor(0, 100, 0)
    solution_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Example
    example_box = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(11.4), Inches(2))
    example_frame = example_box.text_frame
    example_frame.word_wrap = True
    example_frame.text = "Exemple: Dégradation Communication"
    example_frame.paragraphs[0].font.size = Pt(20)
    example_frame.paragraphs[0].font.bold = True
    example_frame.paragraphs[0].font.color.rgb = NAVY
    example_frame.paragraphs[0].space_after = Pt(8)
    
    p = example_frame.add_paragraph()
    p.text = "• AI détecte chute SNR: 20dB → 12dB"
    p.font.size = Pt(16)
    p.space_after = Pt(4)
    
    p = example_frame.add_paragraph()
    p.text = "• Bascule automatique vers antenne de secours"
    p.font.size = Pt(16)
    p.space_after = Pt(4)
    
    p = example_frame.add_paragraph()
    p.text = "• Augmente puissance transmission +10%"
    p.font.size = Pt(16)
    p.space_after = Pt(4)
    
    p = example_frame.add_paragraph()
    p.text = "• Génère alerte critique avec raisonnement"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLUE
    
    # ============================================
    # SLIDE 7: Scenario 2 - Agricultural Detection
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Scénario 2: Détection Agricole & Stress Hydrique"
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = NAVY
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Impact metrics
    metrics = [
        ("99.8%", "Réduction Transmission", "50GB/jour → 100MB/jour"),
        ("150K-1.8M$", "Économies/An", "Par satellite"),
        ("Même Jour", "Alertes vs", "Traitement lendemain")
    ]
    
    for i, (big_num, label, detail) in enumerate(metrics):
        x_pos = 0.8 + i * 4.2
        
        # Metric box
        metric_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                            Inches(x_pos), Inches(1.5),
                                            Inches(3.8), Inches(2))
        metric_box.fill.solid()
        metric_box.fill.fore_color.rgb = LIGHT_BLUE
        metric_box.line.color.rgb = BLUE
        metric_box.line.width = Pt(3)
        
        # Big number
        num_text = slide.shapes.add_textbox(Inches(x_pos + 0.2), Inches(1.7), 
                                           Inches(3.4), Inches(0.8))
        num_frame = num_text.text_frame
        num_frame.text = big_num
        num_frame.paragraphs[0].font.size = Pt(32)
        num_frame.paragraphs[0].font.bold = True
        num_frame.paragraphs[0].font.color.rgb = GOLD
        num_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Label
        label_text = slide.shapes.add_textbox(Inches(x_pos + 0.2), Inches(2.5), 
                                            Inches(3.4), Inches(0.4))
        label_frame = label_text.text_frame
        label_frame.text = label
        label_frame.paragraphs[0].font.size = Pt(14)
        label_frame.paragraphs[0].font.bold = True
        label_frame.paragraphs[0].font.color.rgb = NAVY
        label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Detail
        detail_text = slide.shapes.add_textbox(Inches(x_pos + 0.2), Inches(2.9), 
                                              Inches(3.4), Inches(0.5))
        detail_frame = detail_text.text_frame
        detail_frame.text = detail
        detail_frame.paragraphs[0].font.size = Pt(12)
        detail_frame.paragraphs[0].font.color.rgb = BLACK
        detail_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Process flow
    process_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.4), Inches(2.5))
    process_frame = process_box.text_frame
    process_frame.word_wrap = True
    process_frame.text = "Processus:"
    process_frame.paragraphs[0].font.size = Pt(18)
    process_frame.paragraphs[0].font.bold = True
    process_frame.paragraphs[0].font.color.rgb = NAVY
    process_frame.paragraphs[0].space_after = Pt(8)
    
    p = process_frame.add_paragraph()
    p.text = "Capture Image (500MB-2GB) → Traitement IA → Rapport (1-5MB) → Transmission"
    p.font.size = Pt(16)
    p.font.color.rgb = BLUE
    
    # ============================================
    # SLIDE 8: Technical Feasibility
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Faisabilité Technique"
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = NAVY
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Two columns
    # Left: OS Compatibility
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.5), Inches(5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    left_frame.text = "Compatibilité OS"
    left_frame.paragraphs[0].font.size = Pt(22)
    left_frame.paragraphs[0].font.bold = True
    left_frame.paragraphs[0].font.color.rgb = NAVY
    left_frame.paragraphs[0].space_after = Pt(12)
    
    p = left_frame.add_paragraph()
    p.text = "✓ Linux Embarqué (Yocto)"
    p.font.size = Pt(16)
    p.space_after = Pt(6)
    
    p = left_frame.add_paragraph()
    p.text = "  • Docker/Podman support"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = left_frame.add_paragraph()
    p.text = "✓ VxWorks RTOS"
    p.font.size = Pt(16)
    p.space_after = Pt(6)
    
    p = left_frame.add_paragraph()
    p.text = "  • Conteneurs OCI (depuis 2021)"
    p.font.size = Pt(14)
    p.level = 1
    p.space_after = Pt(4)
    
    p = left_frame.add_paragraph()
    p.text = "  • TensorFlow Lite supporté"
    p.font.size = Pt(14)
    p.level = 1
    
    # Right: Performance
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.5), Inches(5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    right_frame.text = "Performance"
    right_frame.paragraphs[0].font.size = Pt(22)
    right_frame.paragraphs[0].font.bold = True
    right_frame.paragraphs[0].font.color.rgb = NAVY
    right_frame.paragraphs[0].space_after = Pt(12)
    
    p = right_frame.add_paragraph()
    p.text = "Latence: 300-1000ms"
    p.font.size = Pt(16)
    p.space_after = Pt(6)
    
    p = right_frame.add_paragraph()
    p.text = "Mémoire: 2-2.5GB RAM"
    p.font.size = Pt(16)
    p.space_after = Pt(6)
    
    p = right_frame.add_paragraph()
    p.text = "Modèle: <600MB"
    p.font.size = Pt(16)
    p.space_after = Pt(6)
    
    p = right_frame.add_paragraph()
    p.text = "Puissance: <5W"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = BLUE
    
    # ============================================
    # SLIDE 9: Value Proposition by Satellite Type
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Valeur pour Chaque Type de Satellite"
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = NAVY
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Three columns
    types = [
        ("Observation\n(Copernicus, IRIDE)", 
         "• Réduction 80-99% données\n• Alertes temps réel\n• 150K-1.8M$/an économies",
         LIGHT_BLUE),
        ("Communication\n(JSAT-32, SES-17)",
         "• Optimisation réseau\n• Détection interférences\n• -70-80% intervention sol",
         BLUE),
        ("Défense\n(SICRAL, Athena)",
         "• Détection menaces\n• Opération autonome\n• Résilience renforcée",
         NAVY)
    ]
    
    for i, (title, benefits, color) in enumerate(types):
        x_pos = 0.5 + i * 4.3
        
        # Type box
        type_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         Inches(x_pos), Inches(1.5),
                                         Inches(4), Inches(4.5))
        type_box.fill.solid()
        type_box.fill.fore_color.rgb = color
        type_box.line.color.rgb = BLUE
        type_box.line.width = Pt(3)
        
        # Title
        title_text = slide.shapes.add_textbox(Inches(x_pos + 0.2), Inches(1.7), 
                                             Inches(3.6), Inches(0.8))
        title_frame = title_text.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(18)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = WHITE if color == NAVY else NAVY
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Benefits
        benefits_text = slide.shapes.add_textbox(Inches(x_pos + 0.3), Inches(2.7), 
                                                Inches(3.4), Inches(3))
        benefits_frame = benefits_text.text_frame
        benefits_frame.word_wrap = True
        benefits_frame.text = benefits
        benefits_frame.paragraphs[0].font.size = Pt(14)
        benefits_frame.paragraphs[0].font.color.rgb = WHITE if color == NAVY else BLACK
        benefits_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    
    # ============================================
    # SLIDE 10: Implementation Roadmap
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Feuille de Route"
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = NAVY
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Timeline
    phases = [
        ("Phase 1", "Architecture & Faisabilité", "✓ Terminée", GRAY),
        ("Phase 2", "Prototype", "2-3 mois", BLUE),
        ("Phase 3", "Intégration", "3-6 mois", NAVY)
    ]
    
    for i, (phase, desc, timeline, color) in enumerate(phases):
        x_pos = 1 + i * 4
        
        # Phase box
        phase_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                          Inches(x_pos), Inches(1.8),
                                          Inches(3.5), Inches(3.5))
        phase_box.fill.solid()
        phase_box.fill.fore_color.rgb = color
        phase_box.line.color.rgb = BLUE
        phase_box.line.width = Pt(3)
        
        # Phase number
        phase_text = slide.shapes.add_textbox(Inches(x_pos + 0.2), Inches(2), 
                                             Inches(3.1), Inches(0.6))
        phase_frame = phase_text.text_frame
        phase_frame.text = phase
        phase_frame.paragraphs[0].font.size = Pt(24)
        phase_frame.paragraphs[0].font.bold = True
        phase_frame.paragraphs[0].font.color.rgb = WHITE
        phase_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Description
        desc_text = slide.shapes.add_textbox(Inches(x_pos + 0.2), Inches(2.8), 
                                            Inches(3.1), Inches(1.5))
        desc_frame = desc_text.text_frame
        desc_frame.word_wrap = True
        desc_frame.text = desc
        desc_frame.paragraphs[0].font.size = Pt(16)
        desc_frame.paragraphs[0].font.color.rgb = WHITE
        desc_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        
        # Timeline
        time_text = slide.shapes.add_textbox(Inches(x_pos + 0.2), Inches(4.5), 
                                            Inches(3.1), Inches(0.6))
        time_frame = time_text.text_frame
        time_frame.text = timeline
        time_frame.paragraphs[0].font.size = Pt(18)
        time_frame.paragraphs[0].font.bold = True
        time_frame.paragraphs[0].font.color.rgb = GOLD
        time_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Total timeline
    total_box = slide.shapes.add_textbox(Inches(1), Inches(5.8), Inches(11.4), Inches(0.8))
    total_frame = total_box.text_frame
    total_frame.text = "Timeline Total: 5-9 mois jusqu'à solution production-ready"
    total_frame.paragraphs[0].font.size = Pt(18)
    total_frame.paragraphs[0].font.bold = True
    total_frame.paragraphs[0].font.color.rgb = NAVY
    total_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # ============================================
    # SLIDE 11: Call to Action
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "Prêt à Transformer\nvos Systèmes Satellitaires ?"
    title_frame.paragraphs[0].font.size = Pt(48)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = NAVY
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Next steps
    steps_box = slide.shapes.add_textbox(Inches(2), Inches(3.5), Inches(9.4), Inches(2.5))
    steps_frame = steps_box.text_frame
    steps_frame.word_wrap = True
    steps_frame.text = "Prochaines Étapes"
    steps_frame.paragraphs[0].font.size = Pt(24)
    steps_frame.paragraphs[0].font.bold = True
    steps_frame.paragraphs[0].font.color.rgb = BLUE
    steps_frame.paragraphs[0].space_after = Pt(12)
    
    p = steps_frame.add_paragraph()
    p.text = "1. Discussion technique avec équipes Thales"
    p.font.size = Pt(18)
    p.space_after = Pt(8)
    
    p = steps_frame.add_paragraph()
    p.text = "2. Priorisation des cas d'usage"
    p.font.size = Pt(18)
    p.space_after = Pt(8)
    
    p = steps_frame.add_paragraph()
    p.text = "3. Projet pilote (2-3 mois)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = NAVY
    
    # ============================================
    # SLIDE 12: Contact
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(12), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Contact"
    title_frame.paragraphs[0].font.size = Pt(48)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = NAVY
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Contact info
    contact_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9.4), Inches(3.5))
    contact_frame = contact_box.text_frame
    contact_frame.word_wrap = True
    contact_frame.text = "Alexandre GON"
    contact_frame.paragraphs[0].font.size = Pt(32)
    contact_frame.paragraphs[0].font.bold = True
    contact_frame.paragraphs[0].font.color.rgb = NAVY
    contact_frame.paragraphs[0].space_after = Pt(16)
    
    p = contact_frame.add_paragraph()
    p.text = "Architecture & Solutions IA"
    p.font.size = Pt(20)
    p.font.color.rgb = BLUE
    p.space_after = Pt(20)
    
    p = contact_frame.add_paragraph()
    p.text = "alex.gon26@gmail.com"
    p.font.size = Pt(18)
    p.space_after = Pt(8)
    
    p = contact_frame.add_paragraph()
    p.text = "alex.gon@eliago.com"
    p.font.size = Pt(18)
    p.space_after = Pt(8)
    
    p = contact_frame.add_paragraph()
    p.text = "+33 6 27 79 37 72"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = NAVY
    
    # Footer
    footer_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(11.4), Inches(0.8))
    footer_frame = footer_box.text_frame
    footer_frame.text = "Solution techniquement faisable • Économiquement viable • Alignée avec vos objectifs"
    footer_frame.paragraphs[0].font.size = Pt(14)
    footer_frame.paragraphs[0].font.color.rgb = GRAY
    footer_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return prs

if __name__ == "__main__":
    print("Génération de la présentation Thales Alenia Space...")
    presentation = create_presentation()
    output_file = "presentation_thales_alenia.pptx"
    presentation.save(output_file)
    print(f"✓ Présentation créée: {output_file}")
    print(f"✓ Nombre de slides: {len(presentation.slides)}")

